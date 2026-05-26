"""
Generates docs/project_report.pptx — a technical project report for Passive Portfolio Lab.
Run from repo root: python scripts/generate_report_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0d, 0x1f, 0x3c)
BLUE   = RGBColor(0x41, 0x82, 0xb9)
SKY    = RGBColor(0x72, 0xa8, 0xd4)
WHITE  = RGBColor(0xff, 0xff, 0xff)
LIGHT  = RGBColor(0xf7, 0xf9, 0xfc)
TEXT   = RGBColor(0x1a, 0x25, 0x35)
MUTED  = RGBColor(0x6b, 0x72, 0x7e)
GREEN  = RGBColor(0x16, 0xa3, 0x4a)
ACCENT = RGBColor(0xf5, 0x9e, 0x0b)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_para(tf, text, size=14, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p


def section_header(slide, title, subtitle=None, dark=True):
    """Navy top bar with title + optional subtitle."""
    add_rect(slide, 0, 0, 13.33, 1.55, NAVY)
    add_text(slide, title,
             left=0.6, top=0.22, width=12.0, height=0.8,
             size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 left=0.6, top=0.95, width=12.0, height=0.5,
                 size=15, color=SKY)


def kicker(slide, text, left=0.6, top=1.75):
    add_text(slide, text.upper(),
             left=left, top=top, width=6, height=0.35,
             size=10, bold=True, color=BLUE)


def bullet_box(slide, items, left, top, width, height,
               title=None, title_color=NAVY, item_size=13, title_size=14):
    """A labelled bullet list."""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = title
        run.font.size  = Pt(title_size)
        run.font.bold  = True
        run.font.color.rgb = title_color
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.level = 1
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size  = Pt(item_size)
        run.font.color.rgb = TEXT


def add_table(slide, headers, rows, left, top, width, height,
              header_bg=NAVY, header_fg=WHITE, row_size=11, header_size=12):
    cols = len(headers)
    table = slide.shapes.add_table(
        len(rows) + 1, cols,
        Inches(left), Inches(top), Inches(width), Inches(height)
    ).table
    col_w = Inches(width / cols)
    for i in range(cols):
        table.columns[i].width = col_w

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.size  = Pt(header_size)
        run.font.bold  = True
        run.font.color.rgb = header_fg

    # Data rows
    for ri, row in enumerate(rows):
        bg = LIGHT if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            run.font.size  = Pt(row_size)
            run.font.color.rgb = TEXT

    return table


# ── Slides ────────────────────────────────────────────────────────────────────

def slide_cover(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    # Decorative accent bar
    add_rect(s, 0, 3.0, 13.33, 0.06, BLUE)
    # Title
    add_text(s, "被動投資組合實驗室",
             left=1.0, top=1.2, width=11.0, height=1.1,
             size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Passive Portfolio Lab",
             left=1.0, top=2.2, width=11.0, height=0.7,
             size=22, color=SKY, align=PP_ALIGN.CENTER)
    add_text(s, "專案技術報告",
             left=1.0, top=3.3, width=11.0, height=0.6,
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "2026",
             left=1.0, top=4.0, width=11.0, height=0.5,
             size=16, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, "github.com/amandatseng0320/passive-portfolio-lab",
             left=1.0, top=6.6, width=11.0, height=0.4,
             size=11, color=MUTED, align=PP_ALIGN.CENTER)


def slide_overview(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    section_header(s, "專案總覽", "三個平台 · 37 個標的 · 每日自動更新")

    add_text(s, "目標",
             left=0.6, top=1.8, width=3, height=0.4,
             size=14, bold=True, color=NAVY)
    add_text(s,
             "以真實歷史數據回答台灣投資人最常見的投資問題——\n"
             "回測、FIRE 試算、相關性分析、資產篩選，\n"
             "整合為三個免費可用的工具平台。",
             left=0.6, top=2.2, width=6.0, height=1.5,
             size=14, color=TEXT)

    stats = [
        ("37", "精選標的\n(14 台股 ETF ＋ 15 美股 ETF ＋ 8 加密貨幣)"),
        ("3",  "部署平台\n(Streamlit、GitHub Pages、Looker Studio)"),
        ("98+","自動化測試\n(5 模組，覆蓋所有金融公式)"),
        ("24h","資料更新頻率\n(GitHub Actions 每日凌晨一點自動執行)"),
    ]
    for i, (num, label) in enumerate(stats):
        x = 7.0 + (i % 2) * 3.1
        y = 2.0 + (i // 2) * 2.2
        add_rect(s, x, y, 2.9, 1.9, LIGHT)
        add_text(s, num,
                 left=x+0.15, top=y+0.15, width=2.6, height=0.7,
                 size=34, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, label,
                 left=x+0.1, top=y+0.85, width=2.7, height=0.9,
                 size=10, color=MUTED, align=PP_ALIGN.CENTER)


def slide_architecture(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    section_header(s, "系統架構", "四層架構：收集 → 儲存 → 處理 → 呈現")

    layers = [
        ("① 資料收集", ["Yahoo Finance (v8 REST API)", "FRED API（美國 CPI 通膨）", "TWD/USD 匯率（Yahoo TWD=X）"], BLUE),
        ("② 資料儲存", ["Google BigQuery", "raw_prices / asset_metrics", "portfolio_* 預計算結果"], NAVY),
        ("③ 資料處理", ["metrics.py — CAGR、Sharpe、Max DD", "backtest.py — TWD 回測引擎", "fire_calculator.py — FIRE 試算"], BLUE),
        ("④ 呈現平台", ["Streamlit 互動儀表板（即時）", "GitHub Pages 靜態網頁（每日）", "Looker Studio BI 報表（即時）"], NAVY),
    ]

    for i, (title, items, color) in enumerate(layers):
        x = 0.4 + i * 3.2
        add_rect(s, x, 1.8, 3.0, 0.55, color)
        add_text(s, title, left=x+0.1, top=1.85, width=2.8, height=0.45,
                 size=13, bold=True, color=WHITE)
        for j, item in enumerate(items):
            add_text(s, f"• {item}",
                     left=x+0.1, top=2.5+j*0.75, width=2.9, height=0.7,
                     size=11, color=TEXT)

        # Arrow between layers
        if i < 3:
            add_text(s, "→", left=x+3.0, top=2.2, width=0.2, height=0.4,
                     size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    add_text(s, "全流程由 GitHub Actions 排程驅動，每日凌晨一點（台北時間）自動執行。",
             left=0.6, top=6.5, width=12.0, height=0.5,
             size=12, color=MUTED, italic=True)


def slide_pipeline(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    section_header(s, "資料流水線", "GitHub Actions cron → 資料收集 → BigQuery → 三平台部署")

    steps = [
        ("01:00 TW\nGitHub Actions 觸發",   NAVY),
        ("fetch_prices.py\nfetch_macro.py\n37 資產收盤價 + CPI",  BLUE),
        ("metrics.py\nCAGR、Sharpe、\nMax DD 等指標計算",  NAVY),
        ("export_web_data.py\nvalidate_export.py\n寫入 ppl-data.js",  BLUE),
        ("GitHub Pages 部署\nStreamlit 即時讀取\nLooker Studio 即時讀取",  NAVY),
    ]

    box_w = 2.3
    gap   = 0.25
    start_x = 0.3
    for i, (label, color) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        add_rect(s, x, 2.0, box_w, 2.5, color)
        add_text(s, label,
                 left=x+0.1, top=2.15, width=box_w-0.2, height=2.2,
                 size=11, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            ax = x + box_w + 0.02
            add_text(s, "→", left=ax, top=2.9, width=gap+0.1, height=0.5,
                     size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    add_text(s, "BigQuery 作為單一資料來源 (single source of truth)，Streamlit 與 Looker 即時查詢，GitHub Pages 每日匯出靜態快照。",
             left=0.6, top=5.0, width=12.0, height=0.8,
             size=12, color=MUTED, italic=True)


def slide_techstack(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    section_header(s, "科技架構", "Python 3.11 為核心，串接雲端服務與前端平台")

    headers = ["層次", "技術 / 服務", "用途"]
    rows = [
        ("前端介面",   "Streamlit 1.35+, Streamlit-AgGrid", "互動式儀表板"),
        ("視覺化",     "Plotly 5.18+, Chart.js",             "圖表、熱力圖、雷達圖"),
        ("資料處理",   "pandas 2.0+, NumPy 1.26+",           "金融計算與資料框操作"),
        ("資料收集",   "Yahoo Finance v8 REST, FRED API",    "收盤價、匯率、通膨資料"),
        ("資料儲存",   "Google BigQuery",                     "資產指標與每日價格歷史"),
        ("AI 分析",    "Google Gemini 2.5 Flash",             "投資組合策略自然語言摘要"),
        ("靜態網頁",   "HTML / CSS / JavaScript",             "GitHub Pages 無後端部署"),
        ("BI 報表",    "Looker Studio + BigQuery Semantic Layer", "六個預建視圖、可分享報表"),
        ("CI/CD",      "GitHub Actions",                      "每日資料刷新 + 驗證 + 部署"),
        ("測試框架",   "pytest 7.0+, pytest-mock 3.10+",      "單元測試，所有外部 I/O mock"),
        ("語言",       "Python 3.11（主）",                   "2,683 LOC（src + tests）"),
    ]
    add_table(s, headers, rows,
              left=0.5, top=1.7, width=12.3, height=5.4,
              header_size=12, row_size=10)


def slide_assets(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    section_header(s, "37 個精選標的", "台股 ETF ＋ 美股 ETF ＋ 加密貨幣")

    # Three category boxes
    cats = [
        ("台股 ETF（14 檔）",
         ["0050.TW — 台灣 50", "0056.TW — 高股息", "00878.TW — 國泰永續高股息",
          "00919.TW — 群益台灣精選高息", "006208.TW — 富邦台 50",
          "00646.TW — 元大 S&P500", "… 等 14 檔"],
         BLUE),
        ("美股 ETF（15 檔）",
         ["VOO / IVV / SPY — S&P 500", "VTI — 全美市場", "QQQ — 那斯達克 100",
          "VEA — 已開發國家", "VWO — 新興市場", "BND — 美國債券",
          "GLD — 黃金", "VNQ — REITs", "… 等 15 檔"],
         NAVY),
        ("加密貨幣（8 個）",
         ["BTC-USD — 比特幣", "ETH-USD — 以太坊", "BNB-USD — BNB",
          "XRP-USD — 瑞波幣", "SOL-USD — Solana",
          "排除穩定幣與低流動性長尾幣種"],
         BLUE),
    ]
    for i, (title, items, color) in enumerate(cats):
        x = 0.4 + i * 4.3
        add_rect(s, x, 1.75, 4.0, 0.5, color)
        add_text(s, title, left=x+0.1, top=1.8, width=3.8, height=0.4,
                 size=13, bold=True, color=WHITE)
        for j, item in enumerate(items):
            add_text(s, f"• {item}", left=x+0.15, top=2.35+j*0.55,
                     width=3.8, height=0.5, size=11, color=TEXT)

    add_text(s, "篩選標準：流動性充足 ｜ 費用率低 ｜ 覆蓋不重疊（每細分市場 1–2 檔）｜ 至少 3 年歷史資料",
             left=0.6, top=6.5, width=12.0, height=0.5,
             size=12, color=MUTED, italic=True)


def slide_formulas_core(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    section_header(s, "金融公式 — 核心指標", "所有公式均有單元測試覆蓋，無黑盒")

    headers = ["指標", "公式", "說明"]
    rows = [
        ("CAGR 年化報酬",
         "(end_price / start_price)^(1/years) − 1",
         "以完整年數計算複合年增長率"),
        ("波動率 (Volatility)",
         "std(daily_returns) × √252（股票）\nstd(daily_returns) × √365（加密）",
         "年化標準差，加密貨幣採 365 交易日"),
        ("最大回撤 (Max Drawdown)",
         "min((close − rolling_max) / rolling_max)",
         "從歷史最高點的最大跌幅"),
        ("Sharpe 比率",
         "(CAGR − 2.5%) / Volatility",
         "無風險利率預設 2.5%（可調）"),
        ("最差年度 (Worst Year)",
         "min(annual_returns)",
         "歷年中報酬最低的單一自然年"),
    ]
    add_table(s, headers, rows,
              left=0.5, top=1.7, width=12.3, height=5.0,
              header_size=12, row_size=11)


def slide_formulas_backtest(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    section_header(s, "金融公式 — 回測引擎", "以新台幣（TWD）為計算基準的混合回測模型")

    add_text(s, "核心設計",
             left=0.6, top=1.75, width=4, height=0.4,
             size=14, bold=True, color=NAVY)

    points = [
        "Combined TWD Model：第一天一次性投入 ＋ 每月 1 日定期定額",
        "FX 轉換：美股 ETF 以每日實際 TWD/USD 匯率（Yahoo TWD=X）換算",
        "組合價值：所有持倉每日加總；月初注入新資金並重新計算",
        "每月第一個交易日判斷：以該月實際開盤首日注入定期定額",
        "回撤事件標記：識別 12+ 個歷史市場事件（2008 GFC、2020 COVID、\n  2022 升息循環、2025 關稅衝擊等）",
        "恢復期（Recovery）：從谷底回到前高的天數",
    ]
    for i, p in enumerate(points):
        add_text(s, f"• {p}",
                 left=0.7, top=2.2+i*0.72, width=11.8, height=0.65,
                 size=12, color=TEXT)

    add_text(s, "實作檔案：streamlit_dashboard/src/processing/backtest.py",
             left=0.6, top=6.55, width=12.0, height=0.4,
             size=10, color=MUTED, italic=True)


def slide_formulas_fire(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    section_header(s, "金融公式 — FIRE 退休試算", "Financial Independence, Retire Early")

    add_rect(s, 0.5, 1.75, 5.8, 4.8, LIGHT)
    add_text(s, "計算邏輯",
             left=0.7, top=1.9, width=5.4, height=0.4,
             size=14, bold=True, color=NAVY)
    fire_steps = [
        "目標金額 = 年支出 ÷ 提領率\n（預設 4% 法則，可自訂）",
        "每年資產 = 前期資產 × (1 + 加權 CAGR)\n             ＋ 年投入金額",
        "加權 CAGR = Σ (持倉比例ᵢ × CAGRᵢ)",
        "通膨調整：實質報酬 = 名目 CAGR − CPI",
        "上限 50 年；若無法在 50 年內達成則標示未達成",
    ]
    for i, step in enumerate(fire_steps):
        add_text(s, f"• {step}",
                 left=0.7, top=2.4+i*0.85, width=5.4, height=0.8,
                 size=11, color=TEXT)

    add_rect(s, 6.8, 1.75, 5.8, 4.8, NAVY)
    add_text(s, "輸入參數",
             left=7.0, top=1.9, width=5.4, height=0.4,
             size=14, bold=True, color=WHITE)
    inputs = [
        "現有存款（TWD）",
        "每月投入金額（TWD）",
        "目標退休年支出（TWD）",
        "自訂提領率（預設 4%）",
        "通膨率（來自 FRED CPI 或手動輸入）",
        "投資組合配置比例",
    ]
    for i, inp in enumerate(inputs):
        add_text(s, f"• {inp}",
                 left=7.0, top=2.4+i*0.7, width=5.4, height=0.6,
                 size=11, color=WHITE)

    add_text(s, "實作檔案：streamlit_dashboard/src/processing/fire_calculator.py",
             left=0.6, top=6.55, width=12.0, height=0.4,
             size=10, color=MUTED, italic=True)


def slide_tests(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    section_header(s, "測試策略", "5 個模組 · 1,128 行測試程式碼 · 所有外部 I/O mock")

    headers = ["測試檔案", "LOC", "涵蓋範圍"]
    rows = [
        ("test_metrics.py",          "214", "CAGR、波動率（股票 vs 加密）、Max DD、Sharpe、邊界條件"),
        ("test_backtest.py",         "253", "TWD 混合回測引擎、代碼白名單驗證、月投入注入邏輯"),
        ("test_fire_calculator.py",  "201", "FIRE 目標試算、加權 CAGR、通膨調整、50 年上限"),
        ("test_drawdown_events.py",  "230", "回撤事件偵測、12+ 市場事件標記、恢復期計算"),
        ("test_export_web_data.py",  "230", "PPL_ASSETS schema 驗證、PPL_PRICE_HISTORY 結構、TWD 轉換、冪等性"),
    ]
    add_table(s, headers, rows,
              left=0.5, top=1.75, width=12.3, height=3.4,
              header_size=12, row_size=11)

    add_text(s, "測試原則",
             left=0.6, top=5.35, width=4, height=0.4,
             size=14, bold=True, color=NAVY)
    principles = [
        "所有網路 I/O 使用 pytest-mock 模擬（不依賴 BigQuery、Yahoo Finance、FRED、Gemini）",
        "共用測試資料以 conftest.py fixtures 提供（minimal_metrics_df、minimal_prices 等）",
        "每個金融公式均有邊界條件與異常輸入測試",
        "執行：pytest tests/（從 repo root）",
    ]
    for i, p in enumerate(principles):
        add_text(s, f"• {p}",
                 left=0.7, top=5.85+i*0.38, width=12.0, height=0.35,
                 size=11, color=TEXT)


def slide_security(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    section_header(s, "安全性設計", "無硬編碼憑證 · 分層存取控制 · CI/CD OIDC 認證")

    headers = ["元件", "認證方式", "位置"]
    rows = [
        ("Google BigQuery",    "Service Account JSON",              ".env → GOOGLE_APPLICATION_CREDENTIALS"),
        ("Streamlit Cloud",    "Streamlit Secrets（加密儲存）",      "st.secrets['gcp_service_account']"),
        ("應用程式密碼（選用）", "Session-based password gate",       "APP_PASSWORD in st.secrets"),
        ("FRED API（選用）",   "API Key",                           "FRED_API_KEY in .env（無則使用預設值）"),
        ("Gemini AI（選用）",  "API Key",                           "GEMINI_API_KEY in .env（無則隱藏功能）"),
        ("GitHub Actions",    "GCP_SA_KEY encrypted secret",       "OIDC — 無需在 CI 中存儲憑證文件"),
    ]
    add_table(s, headers, rows,
              left=0.5, top=1.75, width=12.3, height=3.2,
              header_size=12, row_size=11)

    add_text(s, "安全設計原則",
             left=0.6, top=5.2, width=4, height=0.4,
             size=14, bold=True, color=NAVY)
    principles = [
        "✓  所有憑證透過環境變數或 Streamlit Secrets 注入，程式碼中無任何明文金鑰",
        "✓  .env.example 提供範例結構，實際 .env 與 credentials.json 列於 .gitignore",
        "✓  export_web_data.py + validate_export.py 雙重驗證，確保輸出資料完整性",
        "✓  GitHub Pages 靜態部署：無後端、無資料庫直連、無 API 金鑰暴露於瀏覽器端",
    ]
    for i, p in enumerate(principles):
        add_text(s, p,
                 left=0.7, top=5.7+i*0.43, width=12.0, height=0.4,
                 size=11, color=TEXT)


def slide_cicd(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    section_header(s, "CI/CD 自動化", "GitHub Actions — 每日一次，全自動無人工介入")

    steps = [
        ("cron: 17:00 UTC\n（台北 01:00）", NAVY),
        ("fetch_prices.py\nfetch_macro.py\n抓取最新收盤價", BLUE),
        ("metrics.py\n計算 CAGR、Sharpe\nMax DD 等指標", NAVY),
        ("export_web_data.py\n匯出 ppl-data.js\n（37 資產 × 日線）", BLUE),
        ("validate_export.py\n資料完整性驗證\n失敗則中止部署", NAVY),
        ("git commit\ngit push\nGitHub Pages 部署", BLUE),
    ]
    box_w = 1.9
    gap = 0.2
    sx = 0.35
    for i, (label, color) in enumerate(steps):
        x = sx + i * (box_w + gap)
        add_rect(s, x, 2.0, box_w, 2.6, color)
        add_text(s, label, left=x+0.08, top=2.15,
                 width=box_w-0.16, height=2.3,
                 size=10.5, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            ax = x + box_w + 0.01
            add_text(s, "→", left=ax, top=3.0, width=gap+0.1, height=0.4,
                     size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    add_text(s,
             "工作流程檔案：.github/workflows/update-and-deploy.yml\n"
             "Streamlit 與 Looker Studio 即時讀取 BigQuery，無需另外部署；只有 GitHub Pages 需要 git push 觸發重新建置。",
             left=0.6, top=5.0, width=12.0, height=0.9,
             size=11, color=MUTED, italic=True)


def slide_changelog(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    section_header(s, "版本更新記錄", "精選重要提交（排除每日自動資料更新）")

    headers = ["Commit", "變更摘要", "類型"]
    rows = [
        ("d875865", "重構 Landing Page：以使用者視角取代技術說明，新增被動投資哲學、精選標的介紹、平台比較表、FAQ",     "feat"),
        ("31e4616", "功能區標語更新：強調完整投資流程與數據驅動決策",                                               "copy"),
        ("3a1e819", "Hero CTA 按鈕錨點修正：滾動至平台介紹區段而非直連 Streamlit",                                  "fix"),
        ("78cc877", "Hero 覆蓋層加入 pointer-events:none，修復按鈕無法點擊問題",                                   "fix"),
        ("177126b", "全站中文文案加入明確斷句 <br />，防止非自然換行",                                              "fix"),
        ("65a4998", "消除重複邏輯、提升效能與可維護性（重構）",                                                    "refactor"),
        ("23165b4", "修正 7 個 P1 正確性錯誤（金融計算相關）",                                                    "fix"),
        ("66923b6", "各平台獨特 favicon、統一按鈕樣式、清理文件結構",                                              "feat"),
        ("b240dfd", "修正響應式設計：卡片 Grid 與全寬版面（手機）",                                                "fix"),
        ("8b45094", "設定所有平台預設語言為正體中文（zh-TW）",                                                    "feat"),
        ("41d9202", "新增 Streamlit Cloud 入口點相容性 shim",                                                    "fix"),
        ("7a476b8", "重組 repo 結構（monorepo → 模組化子目錄）",                                                  "chore"),
        ("3751a14", "新增 CLAUDE.md 專案層級 AI 協作指令",                                                       "chore"),
        ("612b903", "大規模重構：降低重複程式碼、提升可維護性",                                                    "refactor"),
    ]
    add_table(s, headers, rows,
              left=0.4, top=1.7, width=12.5, height=5.5,
              header_size=11, row_size=9)


def slide_compare(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    section_header(s, "三大平台功能比較", "依使用情境選擇最適合的工具")

    headers = ["功能", "Streamlit\n互動儀表板", "GitHub Pages\n靜態網頁", "Looker Studio\n報表"]
    rows = [
        ("六步驟完整分析流程",  "✓", "—", "—"),
        ("即時資料連線",        "✓", "—", "✓"),
        ("AI 投資組合摘要",     "✓", "—", "—"),
        ("歷史回測",            "✓", "—", "✓"),
        ("FIRE 退休試算",       "✓", "—", "✓"),
        ("資產瀏覽與排序",      "✓", "✓", "✓"),
        ("手機友善",            "尚可", "✓", "✓"),
        ("最適情境",            "完整建立分析組合", "快速查閱與視覺瀏覽", "分享報告給他人"),
    ]
    add_table(s, headers, rows,
              left=0.8, top=1.75, width=11.7, height=5.2,
              header_size=13, row_size=12)


def slide_closing(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, 0, 3.0, 13.33, 0.06, BLUE)

    add_text(s, "開源 · 透明 · 可追溯",
             left=1.0, top=1.3, width=11.0, height=0.8,
             size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    points = [
        "完整原始碼發布於 GitHub，MIT 授權",
        "所有金融公式均有對應單元測試，計算邏輯逐行可追蹤",
        "無黑盒：每一條曲線背後都有可驗證的數學公式",
        "每日自動更新，資料來源透明（Yahoo Finance + FRED）",
        "歡迎 Fork、Issue、Pull Request",
    ]
    for i, p in enumerate(points):
        add_text(s, f"• {p}",
                 left=2.5, top=2.2+i*0.65, width=8.5, height=0.6,
                 size=14, color=SKY, align=PP_ALIGN.LEFT)

    add_text(s, "github.com/amandatseng0320/passive-portfolio-lab",
             left=1.0, top=6.5, width=11.0, height=0.5,
             size=13, color=MUTED, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_cover(prs)
    slide_overview(prs)
    slide_architecture(prs)
    slide_pipeline(prs)
    slide_techstack(prs)
    slide_assets(prs)
    slide_formulas_core(prs)
    slide_formulas_backtest(prs)
    slide_formulas_fire(prs)
    slide_tests(prs)
    slide_security(prs)
    slide_cicd(prs)
    slide_changelog(prs)
    slide_compare(prs)
    slide_closing(prs)

    out_dir = os.path.join(os.path.dirname(__file__), "..", "docs")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "project_report.pptx")
    prs.save(out_path)
    print(f"Saved: {os.path.abspath(out_path)}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
